"""Criteo 광고 데이터 파이프라인 포트폴리오 PPT 생성 스크립트."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 색상 팔레트 ────────────────────────────────────────────────────
C_NAVY   = RGBColor(0x1A, 0x2E, 0x4A)
C_BLUE   = RGBColor(0x20, 0x6F, 0xB4)
C_BRONZE = RGBColor(0xCD, 0x7F, 0x32)
C_SILVER = RGBColor(0x70, 0x90, 0xA0)
C_GOLD   = RGBColor(0xC9, 0xA0, 0x2C)
C_GREEN  = RGBColor(0x1E, 0x8B, 0x4C)
C_RED    = RGBColor(0xC0, 0x39, 0x2B)
C_ORANGE = RGBColor(0xD4, 0x6A, 0x0A)
C_LIGHT  = RGBColor(0xF4, 0xF6, 0xF9)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK   = RGBColor(0x1C, 0x1C, 0x1C)
C_GRAY   = RGBColor(0x55, 0x55, 0x55)
C_LGRAY  = RGBColor(0xCC, 0xCC, 0xCC)

W = Inches(13.33)
H = Inches(7.5)


# ── 헬퍼 ──────────────────────────────────────────────────────────
def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, fill=None, border=None):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.line.fill.background()
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s


def tb(slide, text, x, y, w, h,
       size=16, bold=False, color=C_DARK,
       align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def title_bar(slide, title, subtitle=None, bg=C_NAVY):
    bar_h = Inches(1.1)
    rect(slide, 0, 0, W, bar_h, fill=bg)
    tb(slide, title,
       Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.6),
       size=26, bold=True, color=C_WHITE)
    if subtitle:
        tb(slide, subtitle,
           Inches(0.4), Inches(0.7), Inches(12.5), Inches(0.35),
           size=13, color=RGBColor(0xBB, 0xCC, 0xDD))


def bullets(slide, items, x, y, w, h, size=14, color=C_DARK):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        indent = item.startswith("   ")
        p.space_before = Pt(2 if indent else 5)
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size - (1.5 if indent else 0))
        run.font.color.rgb = color


def mixed_bullets(slide, items, x, y, w, h, color=C_DARK):
    """items: list of (text, size, bold)"""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for (text, size, bold) in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        indent = text.startswith("   ") or text.startswith("  ")
        p.space_before = Pt(1 if indent else 4)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def warn(slide, items, x, y, w, h, size=12):
    rect(slide, x, y, w, h,
         fill=RGBColor(0xFF, 0xF3, 0xE0), border=C_ORANGE)
    box = slide.shapes.add_textbox(
        x + Inches(0.1), y + Inches(0.05),
        w - Inches(0.2), h - Inches(0.1))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = C_ORANGE


def chip(slide, text, x, y, w, h, bg=C_BLUE, fg=C_WHITE, size=13, bold=True):
    rect(slide, x, y, w, h, fill=bg)
    tb(slide, text, x, y, w, h,
       size=size, bold=bold, color=fg, align=PP_ALIGN.CENTER)


def sec(slide, text, x, y, w, h=Inches(0.38), color=C_BLUE):
    rect(slide, x, y, w, h, fill=color)
    tb(slide, text,
       x + Inches(0.1), y, w - Inches(0.1), h,
       size=13, bold=True, color=C_WHITE)


# ══════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H


# ── Slide 01: 표지 ───────────────────────────────────────────────
sl = blank_slide(prs)
rect(sl, 0, 0, W, H, fill=C_NAVY)
rect(sl, 0, H - Inches(1.0), W, Inches(1.0), fill=C_BLUE)

tb(sl, "Criteo 광고 데이터 파이프라인",
   Inches(0.7), Inches(1.6), Inches(11.9), Inches(1.2),
   size=44, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
tb(sl, "실시간 수집  →  Medallion 아키텍처  →  BI 시각화",
   Inches(0.7), Inches(2.9), Inches(11.9), Inches(0.7),
   size=22, color=RGBColor(0xA8, 0xC8, 0xE8), align=PP_ALIGN.CENTER)

tags = ["Kafka", "Spark", "S3 Iceberg", "Glue Catalog", "Athena", "Airflow", "Superset"]
tw = Inches(1.4)
sx = (W - tw * len(tags)) / 2
for i, tag in enumerate(tags):
    chip(sl, tag, sx + tw * i, Inches(4.1), Inches(1.3), Inches(0.4), size=12)

tb(sl, "2026.06", Inches(0.7), H - Inches(0.75),
   Inches(3), Inches(0.5), size=14, color=C_WHITE)


# ── Slide 02: 프로젝트 개요 ──────────────────────────────────────
sl = blank_slide(prs)
rect(sl, 0, 0, W, H, fill=C_WHITE)
title_bar(sl, "프로젝트 개요 및 목표")

sec(sl, "프로젝트 배경", Inches(0.4), Inches(1.25), Inches(5.9))
bullets(sl, [
    "도메인: 광고 성과 분석 (Criteo Attribution Dataset)",
    "   HuggingFace 오픈소스 — impression / click / conversion 3종 이벤트",
    "   원본 데이터는 한 행에 조인된 형태 → 실시간 분리 스트림으로 재설계",
    "",
    "목표: 실제 광고 데이터 플랫폼의 수집 → 저장 → 집계 → 시각화",
    "   흐름을 처음부터 끝까지 직접 구현하며 이해",
], Inches(0.5), Inches(1.68), Inches(5.7), Inches(3.5), size=14)

sec(sl, "기술 스택", Inches(6.6), Inches(1.25), Inches(6.35))
rows = [
    ("수집",          "Criteo Dataset (HuggingFace Streaming)",  C_NAVY),
    ("메시지큐",      "Apache Kafka — 3토픽 × 3파티션",           C_BLUE),
    ("스트리밍",      "Spark Structured Streaming",               C_BRONZE),
    ("저장소",        "AWS S3 (Parquet + Iceberg)",                C_SILVER),
    ("메타스토어",    "AWS Glue Data Catalog",                     C_GRAY),
    ("쿼리",          "Amazon Athena (서버리스 SQL)",              C_GREEN),
    ("오케스트레이션","Apache Airflow 3.x",                        C_NAVY),
    ("시각화",        "Apache Superset",                           C_BLUE),
]
for i, (layer, desc, col) in enumerate(rows):
    y = Inches(1.68) + Inches(0.58) * i
    chip(sl, layer, Inches(6.6), y, Inches(1.5), Inches(0.42), bg=col, size=11)
    tb(sl, desc, Inches(8.2), y, Inches(4.7), Inches(0.42), size=12)
rect(sl, Inches(6.6), Inches(1.68), Inches(6.35),
     Inches(0.58) * len(rows), border=C_LGRAY)


# ── Slide 03: 광고 이벤트 이해 (NEW) ────────────────────────────
sl = blank_slide(prs)
rect(sl, 0, 0, W, H, fill=C_WHITE)
title_bar(sl, "광고 이벤트 데이터 이해",
          "impression / click / conversion — 지연 도착과 Attribution Window")

# 이벤트 3종 설명
for i, (name, color, desc, delay) in enumerate([
    ("impression", C_NAVY,
     "광고 노출 이벤트\n사용자 화면에 광고가 보임\ncost(노출 비용) 발생",
     "즉시 발생"),
    ("click", C_BLUE,
     "광고 클릭 이벤트\n사용자가 광고를 클릭\nimpression 발생 후 수초~수일 뒤 도착",
     "impression 후\n수초 ~ 수일"),
    ("conversion", C_GREEN,
     "전환 이벤트\n사용자가 구매/가입 완료\nimpression 발생 후 수일~수십일 뒤 도착",
     "impression 후\n수일 ~ 수십일"),
]):
    bx = Inches(0.3) + Inches(4.3) * i
    rect(sl, bx, Inches(1.25), Inches(4.2), Inches(2.8),
         fill=C_LIGHT, border=color)
    chip(sl, name, bx, Inches(1.25), Inches(4.2), Inches(0.45),
         bg=color, size=15)
    tb(sl, desc, bx + Inches(0.15), Inches(1.75),
       Inches(3.9), Inches(1.7), size=13)
    rect(sl, bx, Inches(4.1), Inches(4.2), Inches(0.45),
         fill=color)
    tb(sl, f"도착 시점: {delay}",
       bx + Inches(0.1), Inches(4.1), Inches(4.0), Inches(0.45),
       size=12, bold=True, color=C_WHITE)

# 핵심 문제
sec(sl, "핵심 문제 — 같은 이벤트인데 도착 시점이 다르다",
    Inches(0.3), Inches(4.7), Inches(12.7), color=C_RED)
bullets(sl, [
    "impression은 지금 저장했는데, 그에 대한 click은 7일 뒤, conversion은 30일 뒤에 도착",
    "→ 파이프라인이 매일 배치로 돌더라도 오늘 impression의 click/conversion을 오늘 알 수 없음",
    "→ Silver 레이어에서 이 '지연 도착' 문제를 2-Stage MERGE로 해결",
], Inches(0.4), Inches(5.12), Inches(12.5), Inches(0.85), size=13)

# Attribution Window
sec(sl, "Attribution Window — '이 impression 덕분에'를 인정하는 시간 한도",
    Inches(0.3), Inches(6.1), Inches(12.7), color=C_NAVY)
bullets(sl, [
    "click window:  7일  — impression 후 7일 이내 click만 이 impression에 귀속 (업계 표준, Google·Meta 동일)",
    "conversion window:  30일  — impression 후 30일 이내 conversion만 귀속",
    "→ 이 window가 Silver 파티션 스캔 범위, Iceberg compaction 35일 기준, Gold KPI 정확도의 핵심 파라미터",
], Inches(0.4), Inches(6.52), Inches(12.5), Inches(0.88), size=13)


# ── Slide 04: 전체 아키텍처 ──────────────────────────────────────
sl = blank_slide(prs)
rect(sl, 0, 0, W, H, fill=C_WHITE)
title_bar(sl, "전체 아키텍처",
          "Criteo Dataset → Kafka → Spark → S3 Iceberg → Athena → Superset")

sl.shapes.add_picture(
    "전체 아키텍처.png",
    Inches(6.397), Inches(0.000), Inches(6.933), Inches(7.500),
)

sec(sl, "설계 원칙", Inches(0.3), Inches(3.0), Inches(6.2))
bullets(sl, [
    "Bronze:  append-only 원본 보존, 언제든 재처리 가능",
    "Silver:  중복 제거 + Attribution 조인, MERGE INTO로 멱등 갱신 (일배치)",
    "Gold:    campaign × date KPI 집계, Superset 직접 조회 대상 (일배치)",
    "",
    "Airflow DAG가 Bronze 확인 → Silver → Gold 순서 보장",
    "배치 사이사이 Athena 검증 쿼리 삽입 (한계 있어 100% 보장 아님)",
], Inches(0.3), Inches(3.45), Inches(6.2), Inches(3.8), size=13)


# ── Slide 05: 왜 Iceberg인가 ──────────────────────────────────────
sl = blank_slide(prs)
rect(sl, 0, 0, W, H, fill=C_WHITE)
title_bar(sl, "왜 Apache Iceberg인가", "일반 Parquet 대비 핵심 기능 비교")

for i, (label, col) in enumerate([
    ("항목", C_NAVY),
    ("일반 Parquet", C_GRAY),
    ("Iceberg (이 프로젝트)", C_BLUE),
]):
    chip(sl, label,
         Inches(0.3) + Inches(4.3) * i, Inches(1.25),
         Inches(4.2), Inches(0.45), bg=col, size=13)

rows_i = [
    ("ACID 쓰기 (MERGE INTO)",
     "직접 지원 없음\n파티션 전체 DROP/INSERT 필요",
     "✅ MERGE INTO 네이티브 지원\n멱등 upsert 가능"),
    ("스냅샷 & Time Travel",
     "불가 — 파일 덮어쓰면\n이전 상태 복구 불가",
     "✅ 스냅샷별 이전 상태 조회\nGold snapshot diff 핵심 기반"),
    ("파티션 변경",
     "파티션 컬럼 변경 시\n파일 전체 재작성",
     "✅ 메타데이터만 변경\n데이터 파일 유지"),
    ("파일 수준 통계",
     "파티션 단위만\n파티션 내 세밀한 pruning 불가",
     "✅ 컬럼 min/max 통계\n파티션 내 파일 pruning 가능"),
    ("COW / MOR 전략",
     "없음 — 항상 파일 전체 재작성",
     "✅ 쓰기 패턴에 따라 선택\nSilver=MOR, Gold=COW"),
]
for i, (item, parquet, iceberg) in enumerate(rows_i):
    y = Inches(1.75) + Inches(1.0) * i
    bg = C_LIGHT if i % 2 == 0 else C_WHITE
    rect(sl, Inches(0.3), y, Inches(12.7), Inches(0.97), fill=bg, border=C_LGRAY)
    tb(sl, item,    Inches(0.4),  y + Inches(0.05), Inches(4.0), Inches(0.85), size=12, bold=True)
    tb(sl, parquet, Inches(4.65), y + Inches(0.05), Inches(3.9), Inches(0.85), size=11, color=C_GRAY)
    tb(sl, iceberg, Inches(8.85), y + Inches(0.05), Inches(4.1), Inches(0.85), size=11, color=C_GREEN)


# ── Slide 06: Medallion 개요 ─────────────────────────────────────
sl = blank_slide(prs)
rect(sl, 0, 0, W, H, fill=C_WHITE)
title_bar(sl, "Medallion 아키텍처 — 레이어별 설계 개요")

cols = [
    ("Bronze", C_BRONZE,
     "원본 보존",
     "S3 Parquet\n(append-only)",
     "raw_date / raw_hour\n(ingest_ts 기준)",
     "Kafka 3토픽 → Spark\nStructured Streaming",
     "없음 (원본 파일이 유일한 재처리 수단)"),
    ("Silver", C_SILVER,
     "분석 가능 단위로 정제",
     "Iceberg\nMOR",
     "event_date\n(impression 발생 기준)",
     "중복 제거 + Attribution\nJoin + MERGE INTO\n(일배치)",
     "MERGE ON (event_id, event_date)"),
    ("Gold", C_GOLD,
     "KPI 집계",
     "Iceberg\nCOW",
     "summary_date\n(event_date와 동일)",
     "Silver → campaign×date\n집계 MERGE INTO\n(일배치)",
     "MERGE ON (summary_date, campaign)"),
]
headers = ["역할", "포맷", "파티션 키", "처리 방식", "멱등성 보장"]
cw = Inches(4.1)
cx0 = Inches(0.3)

for ci, (name, color, *values) in enumerate(cols):
    cx = cx0 + cw * ci
    chip(sl, name, cx, Inches(1.25), cw - Inches(0.1), Inches(0.5), bg=color, size=18)
    for ri, (hdr, val) in enumerate(zip(headers, values)):
        y = Inches(1.82) + Inches(1.05) * ri
        rect(sl, cx, y, cw - Inches(0.1), Inches(1.03),
             fill=C_LIGHT if ri % 2 == 0 else C_WHITE, border=C_LGRAY)
        tb(sl, hdr, cx + Inches(0.1), y + Inches(0.02),
           cw - Inches(0.25), Inches(0.28), size=10, bold=True, color=color)
        tb(sl, val, cx + Inches(0.1), y + Inches(0.3),
           cw - Inches(0.25), Inches(0.65), size=12)


# ── Slide 07: Bronze 상세 ──────────────────────────────────────
sl = blank_slide(prs)
rect(sl, 0, 0, W, H, fill=C_WHITE)
title_bar(sl, "Bronze — 실시간 수집 레이어",
          "Kafka 3토픽  →  Spark Structured Streaming job  →  S3 Parquet")

sec(sl, "Kafka 3토픽 × Spark Streaming job 구성", Inches(0.3), Inches(1.25), Inches(6.1), color=C_BRONZE)
bullets(sl, [
    "ad-impressions  →  Streaming job (cores.max=3)",
    "   impression만 cost 필드 있음 — 스키마 달라 토픽 분리 필수",
    "ad-clicks        →  Streaming job (cores.max=1)",
    "ad-conversions   →  Streaming job (cores.max=1)",
    "",
    "토픽별 독립 Streaming job 컨테이너로 분리",
    "   → impression 처리 지연이 click/conversion 수집에 영향 없음",
    "   → Spark Worker 3대 × 3코어 = 9코어 중 Streaming 5코어 사용",
    "   → 나머지 4코어는 Silver/Gold 일배치용",
], Inches(0.4), Inches(1.68), Inches(5.9), Inches(3.3), size=13)

sec(sl, "raw_date 파티션 = ingest_ts 기준 (수집 시각)", Inches(0.3), Inches(5.05), Inches(6.1), color=C_BRONZE)
bullets(sl, [
    "event_time(이벤트 발생 시각) 기준으로 하면?",
    "   오늘 수집된 conversion이 30일 전 파티션에 기록 → 파일 크기 예측 불가",
    "ingest_ts 기준: 오늘 수집된 모든 이벤트가 오늘 파티션에 균등 적재",
    "   → 장애 복구 시 '이 날짜 파티션부터 Silver 재처리' 범위 특정 가능",
    "   Silver에서 event_date(impression 발생 기준)로 재파티셔닝",
    "   → Bronze는 append-only 원본 저장이므로 Iceberg 불필요",
], Inches(0.4), Inches(5.48), Inches(5.9), Inches(1.65), size=12.5)

sec(sl, "Spark Checkpoint — Streaming 장애 복구", Inches(6.6), Inches(1.25), Inches(6.3), color=C_NAVY)
bullets(sl, [
    "S3에 토픽별 독립 checkpoint 저장",
    "   s3a://…/checkpoints/kafka_to_raw/impressions/",
    "   s3a://…/checkpoints/kafka_to_raw/clicks/",
    "   s3a://…/checkpoints/kafka_to_raw/conversions/",
    "",
    "Streaming job 크래시 후 재시작 시",
    "   → 마지막 처리된 Kafka offset에서 이어서 처리",
    "   → 중복/누락 없이 재개 가능",
    "",
    "checkpoint 경로를 토픽별로 분리하는 이유",
    "   공유하면 offset 상태가 섞여 재처리 시 누락 또는 중복 발생",
], Inches(6.7), Inches(1.68), Inches(6.1), Inches(4.0), size=13)

warn(sl, [
    "⚠️ 한계",
    "• Kafka 단일 브로커 (replication factor=1) — 브로커 장애 시 메시지 손실",
    "  → 운영 환경에선 3+ 브로커 클러스터 + replication 구성 필요",
    "• Spark Worker가 단일 머신에 모두 있어 실제 분산 환경 아님 (흉내 수준)",
    "• producer.py 재시작 시 전체 데이터 처음부터 재발행 — 모든 이벤트 중복 가능 (Silver MERGE에서 흡수)",
], Inches(6.6), Inches(5.7), Inches(6.3), Inches(1.65))


# ── Slide 08: Silver 상세 ──────────────────────────────────────
sl = blank_slide(prs)
rect(sl, 0, 0, W, H, fill=C_WHITE)
title_bar(sl, "Silver — Attribution 2-Stage 일배치",
          "Dedup + Attribution Join + Iceberg MERGE INTO (MOR)")

sec(sl, "왜 2-Stage인가 — 이벤트 도착 시점 문제", Inches(0.3), Inches(1.25), Inches(6.1), color=C_SILVER)
bullets(sl, [
    "어제(6/24) Bronze raw_date=6/24에 들어오는 것:",
    "   impression: 어제 발생 → event_date = 6/24",
    "   click:      어제 도착, 원본 impression은 6/17~6/24 (click window 7일)",
    "   conversion: 어제 도착, 원본 impression은 5/25~6/24 (window 30일)",
    "",
    "1-Stage JOIN (어제 impression + 어제 click)으로 하면?",
    "   → click은 어제 Bronze에 있지만 6/17 impression은 어제 Bronze에 없음 → 매칭 불가",
    "   → 최대 30일치 Bronze를 매일 전체 스캔해야 해결 → 너무 비쌈",
    "   (conversion window 30일 기준)",
    "",
    "Stage 1: 어제 Bronze의 impression → Silver INSERT",
    "   event_date = impression 발생 시각 기준 (여러 파티션 분산 가능)",
    "   이미 Silver에 있는 impression은 INSERT 안 함 (멱등)",
    "   click=0, conversion=0 초기값으로 적재",
    "Stage 2: 어제 Bronze의 click/conversion → Silver 7일/30일 파티션에서",
    "   eid로 impression 찾아 MERGE UPDATE",
    "   → Silver가 중간 저장소 역할, 30일치 Bronze 스캔 불필요",
], Inches(0.4), Inches(1.68), Inches(5.9), Inches(4.5), size=12.5)

sec(sl, "왜 MOR(Merge-On-Read)인가", Inches(0.3), Inches(6.2), Inches(6.1), color=C_SILVER)
bullets(sl, [
    "UPDATE 대상(click/conversion 매칭 impression): 전체의 ~2%",
    "COW: 2% 행이 바뀌어도 그 파티션 파일 전체 재작성",
    "MOR: delta file만 append → 쓰기 비용 절감",
    "읽기 시 delta 병합 오버헤드 → 유지보수 DAG의 daily compaction이 해결",
], Inches(0.4), Inches(6.6), Inches(5.9), Inches(0.85), size=12.5)

warn(sl, [
    "⚠️  백필 시 Silver도 명시 재처리 모드(--run-date-start 지정) 사용 필요",
    "   일배치 기본 모드(snapshot diff)는 같은 날 여러 번 실행 시 변경 감지가 깨짐",
], Inches(0.4), Inches(7.55), Inches(5.9), Inches(0.6))

sec(sl, "Attribution 파티션 스캔 범위", Inches(6.6), Inches(1.25), Inches(6.3), color=C_NAVY)
bullets(sl, [
    "Stage 2 Silver 스캔 시 전체 파티션을 열지 않음",
    "",
    "click JOIN 조건:",
    "   ① click_ts - impression_ts 가 0 ~ 7일 사이  (attribution 판정)",
    "   ② Silver event_date >= click 기준 7일 전  (파티션 범위 제한)",
    "   → Iceberg manifest가 event_date별 파티션 정보를 가짐",
    "     → ② 기준 해당 날짜 이후 파티션만 파일 오픈",
    "   → manifest는 파일 단위 컬럼 min/max도 기록",
    "     → eid 값이 파일 min/max 밖이면 그 파일도 추가 스킵 가능",
    "   → 남은 파일 안에서 eid로 impression 행 검색 (row-level index는 없음)",
    "",
    "conversion JOIN 조건:",
    "   마찬가지로 conversion 기준 30일치 파티션만 스캔",
    "",
    "MERGE INTO ON (event_id, event_date)",
    "   Stage 1 — NOT MATCHED → INSERT (신규 impression 초기값 적재)",
    "   Stage 2 — MATCHED AND click=0 → UPDATE SET click=1",
    "             MATCHED AND conversion=0 → UPDATE SET conversion=1",
], Inches(6.7), Inches(1.68), Inches(6.1), Inches(5.1), size=12.5)


# ── Slide 09: Gold 상세 ──────────────────────────────────────
sl = blank_slide(prs)
rect(sl, 0, 0, W, H, fill=C_WHITE)
title_bar(sl, "Gold — Snapshot Diff + KPI 집계",
          "Silver 변경분만 재집계  →  Iceberg MERGE INTO (COW)  →  일배치")

sec(sl, "Snapshot Diff — Silver 변경 파티션만 재집계", Inches(0.3), Inches(1.25), Inches(6.1), color=C_GOLD)
bullets(sl, [
    "아이디어: 일배치 → 오늘 Silver에서 바뀐 것만 재집계하면 충분",
    "   Silver 전체를 매번 재스캔하는 건 낭비",
    "   (attribution window 30일치만 봐도 되지만 그것도 매일 전부 재집계하면 비효율)",
    "",
    "구현: Silver Iceberg 스냅샷으로 변경 파티션 감지",
    "   ① 오늘 Silver 첫 번째 커밋의 parent = 배치 이전 상태",
    "   ② 현재 스냅샷 vs parent 스냅샷을 event_date별 행수 비교",
    "   ③ 달라진 event_date = 오늘 변경 파티션 → Gold 재집계 대상",
    "      (당일 신규 + 지연 attribution으로 변경된 과거 파티션 자동 포함)",
    "",
    "MERGE INTO ON (summary_date, campaign)",
    "   MATCHED  → UPDATE SET *  (재집계 결과로 덮어쓰기)",
    "   NOT MATCHED  → INSERT  (신규 날짜)",
], Inches(0.4), Inches(1.68), Inches(5.9), Inches(4.2), size=12.5)

sec(sl, "왜 COW(Copy-On-Write)인가", Inches(0.3), Inches(5.95), Inches(6.1), color=C_GOLD)
bullets(sl, [
    "Gold는 Superset/Athena에서 가장 많이 읽히는 테이블 → 읽기 성능이 최우선",
    "COW: 파티션 파일 재작성 후 항상 clean → Athena 스캔 성능 일정",
], Inches(0.4), Inches(6.35), Inches(5.9), Inches(1.05), size=12.5)

sec(sl, "Gold KPI — campaign_summary 스키마", Inches(6.6), Inches(1.25), Inches(6.3), color=C_NAVY)
kpis = [
    ("impressions",    "COUNT(*)"),
    ("clicks",         "SUM(click)"),
    ("conversions",    "SUM(conversion)"),
    ("unique_users",   "COUNT(DISTINCT uid)"),
    ("CTR",            "clicks / impressions × 100"),
    ("CVR",            "conversions / clicks × 100  (clicks=0 → NULL)"),
    ("CPC",            "total_cost / clicks  (clicks=0 → NULL)"),
    ("CPA",            "total_cost / conversions  (conversions=0 → NULL)"),
    ("CPM",            "total_cost / impressions × 1,000"),
    ("frequency",      "impressions / unique_users"),
    ("avg_conv_delay", "AVG(conversion_delay_sec ≥ 0)  — sentinel(-1) 제외"),
    ("CTC / VTC",      "click=1&conversion=1  /  click=0&conversion=1"),
]
for i, (kpi, formula) in enumerate(kpis):
    y = Inches(1.75) + Inches(0.42) * i
    bg = C_LIGHT if i % 2 == 0 else C_WHITE
    rect(sl, Inches(6.6), y, Inches(6.3), Inches(0.41), fill=bg, border=C_LGRAY)
    tb(sl, kpi,     Inches(6.7),  y + Inches(0.04), Inches(1.8), Inches(0.35),
       size=11, bold=True, color=C_GOLD)
    tb(sl, formula, Inches(8.55), y + Inches(0.04), Inches(4.3), Inches(0.35),
       size=11)

warn(sl, [
    "⚠️  백필 시 Gold는 명시 재처리 모드(--run-date-start 지정) 사용 필요",
], Inches(6.6), Inches(7.08), Inches(6.3), Inches(0.35))


# ── Slide 10: Airflow DAG ────────────────────────────────────────
sl = blank_slide(prs)
rect(sl, 0, 0, W, H, fill=C_WHITE)
title_bar(sl, "운영 자동화 — Airflow DAG 구조",
          "3개 DAG + 검증 게이트 + ExternalTaskSensor 의존성")

sec(sl, "DAG 의존 관계", Inches(0.3), Inches(1.25), Inches(8.5))
sl.shapes.add_picture(
    "airflow DAG 구조.png",
    Inches(0.300), Inches(1.680), Inches(6.365), Inches(5.241),
)

sec(sl, "검증 게이트 설계", Inches(9.0), Inches(1.25), Inches(4.0))
bullets(sl, [
    "패턴: 0행 = 정상, 1행+ = AirflowException",
    "",
    "※ 오늘 = 배치 실행일 / 어제 = Bronze 처리 대상일",
    "",
    "verify_silver_snapshot",
    "   오늘 배치가 Silver Iceberg 커밋을 남겼는가",
    "   (Spark exit 0이어도 커밋 없으면 감지)",
    "",
    "verify_silver_rows",
    "   어제(처리 대상) Silver event_date 파티션 행수 > 0",
    "",
    "verify_gold_snapshot",
    "   오늘 배치가 Gold Iceberg 커밋을 남겼는가",
    "",
    "verify_silver_gold_consistency",
    "   Silver COUNT(*) vs Gold SUM(impressions)",
    "   차이 10% 이내 + Gold 파티션 존재 확인",
    "   (10%: attribution 지연으로 Silver/Gold 완전 일치 불가",
    "   → 보수적 기준)",
], Inches(9.1), Inches(1.68), Inches(3.9), Inches(4.2), size=12)

warn(sl, [
    "⚠️ 한계",
    "• 검증 범위: 어제 파티션 1개만 확인",
    "  Attribution으로 갱신된 어제 이전 파티션 Silver↔Gold 정합성 미커버",
    "• 이 DAG 구조는 Silver/Gold 테이블이 각 1개라는 전제에 의존",
    "  새 테이블 추가 시 전체 DAG가 무너지는 구조 → 레이어별 유지보수 DAG 분리가 더 적합한 설계였음",
    "• 백필: --run-date-start/end 인자 있으나 Airflow UI에서 날짜 지정 연결 미구현",
    "  — docker compose 명령 직접 실행 필요",
], Inches(0.3), Inches(5.65), Inches(12.7), Inches(1.75))


# ── Slide 11: Iceberg 유지보수 ──────────────────────────────────
sl = blank_slide(prs)
rect(sl, 0, 0, W, H, fill=C_WHITE)
title_bar(sl, "Iceberg 유지보수 전략",
          "Compaction / Expire / Orphan / Manifest 재편성 / metadata.json 보존")

sec(sl, "Silver (MOR) 일간 유지보수", Inches(0.3), Inches(1.25), Inches(6.1), color=C_SILVER)
bullets(sl, [
    "compact_silver  (rewrite_data_files)",
    "   대상: event_date >= 오늘-35일",
    "   MOR MERGE가 매일 delete file 누적 → base + delta 병합 오버헤드 제거",
    "   왜 35일?: attribution window 30일 + 수집 지연 여유 5일",
    "",
    "expire_silver_snapshots",
    "   older_than: 오늘-31일  (31일 초과 스냅샷 제거, 최대 31일 보존)",
    "",
    "remove_silver_orphans",
    "   older_than: 오늘-4일  (진행 중 쓰기와 충돌 방지)",
    "",
    "metadata.json 보존 설정 (TBLPROPERTIES)",
    "   previous-versions-max = 21  (하루 3커밋 × 7일 = 21)",
    "   커밋마다 새 metadata.json 생성 → 무한 누적 방지",
    "   ※ 스냅샷 보존(31일) > metadata.json 버전 보존(7일치)",
    "     스냅샷: 롤백·시간여행 위해 길게 / metadata.json 파일 수: 짧게 관리",
], Inches(0.4), Inches(1.68), Inches(5.9), Inches(4.8), size=12.5)

sec(sl, "Gold (COW) 일간 유지보수", Inches(6.6), Inches(1.25), Inches(6.3), color=C_GOLD)
bullets(sl, [
    "Compaction: 현재 규모에서 불필요",
    "   COW는 delete file 없음 → MOR처럼 매일 필수 아님",
    "   현재 데이터 규모: 파티션당 Parquet 파일 1개 수준 → 소형 파일 누적 없음",
    "   데이터가 커지면 소형 파일 병합을 위한 compaction 필요해질 수 있음",
    "",
    "expire_gold_snapshots",
    "   older_than: 오늘-31일",
    "",
    "remove_gold_orphans",
    "   older_than: 오늘-4일",
    "",
    "metadata.json 보존 설정",
    "   previous-versions-max = 7  (하루 1커밋 × 7일)",
], Inches(6.7), Inches(1.68), Inches(6.1), Inches(3.5), size=12.5)

sec(sl, "월간 유지보수 (매월 1일)", Inches(6.6), Inches(5.25), Inches(6.3), color=C_NAVY)
bullets(sl, [
    "rewrite_manifests (Silver + Gold 병렬)",
    "   스냅샷 누적으로 manifest 파편화 → 소수의 큰 파일로 재편성",
    "   데이터 파일 변경 없이 쿼리 플래닝 오버헤드 감소",
    "   현재 규모(일 파티션 1개)에서는 월 1회로 충분",
], Inches(6.7), Inches(5.68), Inches(6.1), Inches(1.65), size=12.5)

sec(sl, "MERGE와 Compaction 충돌 방지", Inches(0.3), Inches(6.55), Inches(12.7), color=C_NAVY)
bullets(sl, [
    "ExternalTaskSensor: medallion DAG gold_batch SUCCESS 감지 후 maintenance 시작",
    "→ Silver/Gold MERGE 완전히 끝난 뒤 Compaction 실행 → Iceberg 낙관적 잠금 충돌 없음",
], Inches(0.4), Inches(6.93), Inches(12.5), Inches(0.5), size=12.5)


# ── Slide 12: Superset 대시보드 ──────────────────────────────────
sl = blank_slide(prs)
rect(sl, 0, 0, W, H, fill=C_WHITE)
title_bar(sl, "Superset 대시보드", "Gold → 비즈니스 KPI  |  Silver → 운영 모니터링")

sec(sl, "비즈니스 대시보드 — Gold 기반 (Campaign팀/BI팀)", Inches(0.3), Inches(1.25), Inches(8.0), color=C_GOLD)
bullets(sl, [
    "데이터 소스: Gold campaign_summary",
    "(집계 테이블 → Athena 스캔 비용 낮음)",
    "",
    "business/01  30일 캠페인별 KPI 추이",
    "   CTR / CVR / CPC / CPA / CPM / frequency",
    "   click-through conversion vs view-through conversion",
    "",
    "business/02  Conversion 지연 분포 (주간 refresh)",
    "   p25 / p50 / p75 / p90 / p95",
    "   Attribution window(30일) 적정성 검증 지표",
    "",
    "Superset Dataset 캐싱 (미구현 — 적용 시 비용 절감 가능)",
    "   24h 캐싱 설정 시 하루 1번 Athena 스캔으로 여러 팀 조회 커버",
], Inches(0.4), Inches(1.68), Inches(7.8), Inches(3.9), size=13)

sec(sl, "운영 대시보드 — Silver 기반 (DE용)", Inches(0.3), Inches(5.65), Inches(8.0), color=C_NAVY)
bullets(sl, [
    "Silver를 데이터 소스로 사용 — 파이프라인 정상 동작 여부 확인 전용",
    "   일별 볼륨 추이 / attribution 커버리지 / 단계별 파이프라인 지연 추이 (p50)",
    "   운영 이상 신호(급증/급감/attribution 0%)를 시각적으로 모니터링",
], Inches(0.4), Inches(6.08), Inches(7.8), Inches(1.25), size=13)

sl.shapes.add_picture(
    "대시보드.png",
    Inches(5.200), Inches(1.100), Inches(7.958), Inches(3.072),
)
sl.shapes.add_picture(
    "대시보드_ops.png",
    Inches(5.980), Inches(3.373), Inches(5.061), Inches(4.007),
)


# ── Slide 13: 회고 ────────────────────────────────────────────────
sl = blank_slide(prs)
rect(sl, 0, 0, W, H, fill=C_WHITE)
title_bar(sl, "회고")

sec(sl, "설계 한계 — 고려하지 못한 것들", Inches(0.3), Inches(1.25), Inches(6.1), color=C_RED)
mixed_bullets(sl, [
    ("인프라 HA 미고려",                                                            10, True),
    ("   Kafka 단일 브로커 — 장애 시 메시지 손실",                                   9, False),
    ("   Spark Worker가 단일 물리 머신 — 실제 분산 아님",                             9, False),
    ("   장애 복구 절차 미설계",                                                       9, False),
    ("플랫폼 확장성",                                                                10, True),
    ("   Silver/Gold 테이블이 여러 개 생기는 상황을 생각 못 했음",                    9, False),
    ("   지금 구조는 테이블 1:1 하드코딩 — 테이블 늘면 DAG 전체 수정",               9, False),
    ("   레이어별 유지보수 DAG 분리가 더 적합했을 것",                                 9, False),
    ("   Bronze도 Silver 테이블 늘면 읽히는 빈도 증가",                               9, False),
    ("   → Bronze compaction Spark job 별도 필요해질 수 있음",                        9, False),
    ("스케일아웃 전략 없음",                                                          10, True),
    ("   데이터 10x 증가 시 파티션 전략 / 코어 배분 어떻게 바꿔야 하는지 고려 못 했음", 9, False),
    ("날짜 기준 UTC 고정 (UTC vs KST)",                                              10, True),
    ("   Kafka·Spark·Airflow·S3 모두 UTC 기준으로 동작→ 날짜 경계 처리 주의",        9, False),
    ("   현재 Airflow를 KST 낮 시간대에 실행해 회피 중 (근본 해결 아님)",              9, False),
    ("백필 운영 절차 미흡",                                                           10, True),
    ("   CLI 인자는 있는데 Airflow UI에서 날짜 지정하는 연결 미구현",                  9, False),
    ("Compaction 정렬 전략 미고려",                                                   10, True),
    ("   현재 bin-pack: 파일 크기만 맞춤, 정렬 없음",                                  9, False),
    ("   sort/z-order 적용 시 attribution JOIN·Superset 쿼리 성능 개선 가능",          9, False),
    ("Gold 집계 버그 — 발표 준비 중 발견 (미수정)",                                  10, True),
    ("   Silver 스냅샷 비교로 변경 파티션만 Gold 재집계하는 방식인데",                 9, False),
    ("   click·conversion은 기존 행 값만 바뀌어 행 수 그대로 → 감지 누락",            9, False),
], Inches(0.4), Inches(1.68), Inches(5.9), Inches(5.5))

sec(sl, "프로젝트 진행 후기", Inches(6.6), Inches(1.25), Inches(6.3), color=C_NAVY)
mixed_bullets(sl, [
    ("CS 기반 부족 — 환경 설정에서 막혔음",                                           12, True),
    ("   환경변수·포트·네트워크·엔드포인트 등 처음 들어보는 개념들이 많아 이해 어려움", 10, False),
    ("   .config() 자격증명·docker-compose·Dockerfile 설정은",                        10, False),
    ("   Claude Code에 맡기고 '앱이 돌아가면 된 거겠지'로 진행",                       10, False),
    ("기능 단위 테스트를 안 했음",                                                     12, True),
    ("   기능 여러 개 붙이고 파이프라인 한번에 돌려서 오류 역추적",                    10, False),
    ("   → 기능 하나 붙일 때마다 소규모 테스트했으면 더 효율적",                       10, False),
    ("적재 데이터 직접 검증 안 했음",                                                  12, True),
    ("   파이프라인이 오류 없이 돌면 데이터가 맞겠지 하고 넘어감",                    10, False),
    ("  → 테이블 만들 때마다 소규모 Athena 쿼리로 여러 케이스 확인했어야 함",          10, False),
    ("Spark 자원 배분 고민 부족",                                                      12, True),
    ("   worker·executor 조정으로 처리량 늘린다는 건 알지만",                          10, False),
    ("   job별 구체적 자원 설계는 못 해봄",                                            10, False),
    ("Claude Code 처음 써봤는데 AI가 제안한 코드 읽기부터 막히는 문제",               12, True),
    ("   구현 됐다고 하면 코드 안 읽고 패스하는 식으로 진행",                          10, False),
    ("   코드 읽기부터 막히는 수준인데 직접 쓰기까지는 갭이 얼마나 클지",             10, False),
    ("   바이브코딩·노코드 추세인데 코드 읽기/쓰기에 대한 공부가",                    10, False),
    ("   어느 정도로 필요할지 방향이 안 잡힘",                                         10, False),
], Inches(6.7), Inches(1.68), Inches(6.1), Inches(5.5))


# ── Slide 14: Q&A ─────────────────────────────────────────────
sl = blank_slide(prs)
rect(sl, 0, 0, W, H, fill=C_NAVY)
rect(sl, 0, H - Inches(1.2), W, Inches(1.2), fill=C_BLUE)

tb(sl, "감사합니다",
   Inches(0.7), Inches(2.2), Inches(11.9), Inches(1.4),
   size=60, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
tb(sl,
   "Kafka  →  Spark Streaming  →  S3 Iceberg (MOR / COW)  →  Athena  →  Superset\n"
   "Attribution Window  |  Snapshot Diff  |  MERGE INTO 멱등성  |  Airflow 자동화",
   Inches(0.7), Inches(3.8), Inches(11.9), Inches(1.0),
   size=15, color=RGBColor(0xA8, 0xC8, 0xE8), align=PP_ALIGN.CENTER)


# ── 저장 ──────────────────────────────────────────────────────────
out = "/home/sandy/metacode-project/criteo_pipeline_portfolio.pptx"
prs.save(out)
print(f"저장 완료: {out}")
print(f"슬라이드 수: {len(prs.slides)}")
